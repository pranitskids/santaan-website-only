from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph


ROOT = Path(__file__).resolve().parents[1]
SCREENSHOT_DIR = ROOT / "docs" / "manuals" / "screenshots"


BRAND_TEAL = RGBColor(34, 85, 82)
BRAND_AMBER = RGBColor(230, 164, 78)
BRAND_CREAM = RGBColor(246, 242, 235)
BRAND_INK = RGBColor(31, 38, 47)
BRAND_MUTED = RGBColor(101, 109, 119)
BRAND_STONE = RGBColor(224, 220, 213)
WHITE = RGBColor(255, 255, 255)


@dataclass(frozen=True)
class DeckSection:
    title: str
    subtitle: str
    bullets: Sequence[str]
    image_name: str | None = None


@dataclass(frozen=True)
class DeckSpec:
    slug: str
    title: str
    subtitle: str
    kicker: str
    sections: Sequence[DeckSection]

    @property
    def pptx_out(self) -> Path:
        return ROOT / "output" / "pptx" / f"{self.slug}.pptx"

    @property
    def pdf_out(self) -> Path:
        return ROOT / "output" / "pdf" / f"{self.slug}.pdf"


VISUAL_DECK = DeckSpec(
    slug="santaan-crm-visual-onboarding-manual-2026-03-29",
    title="Santaan CRM Visual Onboarding Manual",
    subtitle="A practical onboarding deck built from the live CRM screens",
    kicker="TEAM ONBOARDING",
    sections=[
        DeckSection(
            title="Login And Access",
            subtitle="Every team member enters through the same controlled portal",
            bullets=[
                "Use the role dropdown if available, then enter username or email plus the 6-digit PIN.",
                "Magic link is the backup path only for whitelisted email accounts.",
                "Admins should use User Access to create staff users and reset PINs.",
            ],
            image_name="login-page.png",
        ),
        DeckSection(
            title="Today Tab",
            subtitle="The operating home screen for the day",
            bullets=[
                "Start every day here to read the rider, review standup prompts, and open the right workflow tabs.",
                "Leadership uses this view to jump into Analytics, Spend, Daily Command, and Workboard.",
                "The snapshot at the top helps spot follow-up gaps, stale leads, and integration health quickly.",
            ],
            image_name="dashboard-today.png",
        ),
        DeckSection(
            title="Publish Content",
            subtitle="Direct website publishing from inside the CRM",
            bullets=[
                "Content managers can publish announcements, clinical insights, and doctor updates without Medium.",
                "The form accepts plain paragraphs and turns them into a live website post with a generated slug.",
                "Use this for fast, controlled publishing with no developer handoff.",
            ],
            image_name="publish-content.png",
        ),
        DeckSection(
            title="NeoDove Ops",
            subtitle="The calling sync control room",
            bullets=[
                "Use this screen to catch missing owners, missing follow-ups, stale sync, and webhook issues.",
                "This is the first place managers should look when leads are leaking after a NeoDove event.",
                "The current live NeoDove workflow model is campaign-by-campaign with Lead Created, Call Connected, and Call Not Connected.",
            ],
            image_name="neodove-ops.png",
        ),
        DeckSection(
            title="User Access",
            subtitle="PIN-based onboarding and role control",
            bullets=[
                "Admins can create new staff users, assign roles, reset PINs, and bulk onboard teams from this screen.",
                "Role assignment controls which tabs appear for each user, so this screen is central to usability.",
                "Disable stale access promptly so the dashboard stays clean and secure.",
            ],
            image_name="user-access.png",
        ),
        DeckSection(
            title="Daily Team Rhythm",
            subtitle="What makes the system commercially useful",
            bullets=[
                "CEO: review Today, Command Center, Analytics, Spend, NeoDove Ops, and Workboard.",
                "Content: work from Publish Content, Draft Content, Content Insights, and Workboard.",
                "Telecalling: use Hot Leads, Follow-ups, All Contacts, and NeoDove Ops to close loops the same day.",
                "The CRM becomes valuable only when statuses, notes, owners, and next follow-ups are updated consistently.",
            ],
        ),
    ],
)


CEO_DECK = DeckSpec(
    slug="santaan-crm-ceo-brief-2026-03-29",
    title="Santaan CRM Founder Brief",
    subtitle="What is live now, why it matters, and how to run the next phase with discipline",
    kicker="CEO BRIEF",
    sections=[
        DeckSection(
            title="What Is Live Now",
            subtitle="The CRM has moved from build mode into operating mode",
            bullets=[
                "NeoDove to CRM webhook flow is live and acceptance-tested with successful 200 responses.",
                "WhatsApp AI enrichment is active and can push lead context into the calling loop.",
                "Content publishing is live inside the CRM, so brand and clinic teams can ship without external tools.",
            ],
            image_name="dashboard-today.png",
        ),
        DeckSection(
            title="Where Revenue Leakage Can Now Be Seen",
            subtitle="The key gain is not just automation, it is visibility",
            bullets=[
                "NeoDove Ops surfaces missing owners, missing follow-ups, stale sync, and status drift in one place.",
                "This lets management catch leakage before leads disappear between marketing, telecalling, and counseling.",
                "Operational truth is now closer to real time instead of waiting for verbal updates or spreadsheets.",
            ],
            image_name="neodove-ops.png",
        ),
        DeckSection(
            title="What Teams Can Operate Directly",
            subtitle="The system is useful because each function can act inside its own role screen",
            bullets=[
                "Content can publish to the live website directly from the CRM.",
                "Admins can control onboarding, roles, and PIN resets from User Access.",
                "This reduces tool switching and makes adoption easier for non-technical staff.",
            ],
            image_name="publish-content.png",
        ),
        DeckSection(
            title="The Adoption Rule",
            subtitle="The business benefit now depends more on behavior than on code",
            bullets=[
                "Every active lead must have status, owner, note, and next follow-up.",
                "Leadership should review the same dashboard every day and use it as the source of truth.",
                "If teams stop updating the system, the CRM becomes decorative instead of operational.",
            ],
            image_name="user-access.png",
        ),
    ],
)


DECKS = [VISUAL_DECK, CEO_DECK]


def ensure_dirs() -> None:
    (ROOT / "output" / "pptx").mkdir(parents=True, exist_ok=True)
    (ROOT / "output" / "pdf").mkdir(parents=True, exist_ok=True)


def fill_rect(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, color, bold=False, font_name="Aptos", align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_bullet_block(slide, left, top, width, height, bullets, font_size=16):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = BRAND_STONE

    tf = panel.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.font.size = Pt(font_size)
        p.font.color.rgb = BRAND_INK
        p.font.name = "Aptos"
    return panel


def image_fit(image_path: Path, max_w_in: float, max_h_in: float) -> tuple[float, float]:
    with Image.open(image_path) as img:
        width, height = img.size
    scale = min(max_w_in / width, max_h_in / height)
    return width * scale, height * scale


def add_cover_slide(prs: Presentation, spec: DeckSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_rect(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height), BRAND_TEAL)

    fill_rect(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.6), Inches(2.3), Inches(0.46)), BRAND_AMBER)
    add_textbox(slide, Inches(0.85), Inches(0.68), Inches(2.0), Inches(0.25), spec.kicker, 11, WHITE, bold=True)

    add_textbox(slide, Inches(0.75), Inches(1.35), Inches(7.2), Inches(1.2), spec.title, 28, WHITE, bold=True)
    add_textbox(slide, Inches(0.75), Inches(2.45), Inches(7.4), Inches(0.9), spec.subtitle, 15, RGBColor(225, 232, 230))

    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.35), Inches(5.65), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = BRAND_CREAM
    card.line.color.rgb = BRAND_CREAM
    tf = card.text_frame
    tf.clear()
    summary_lines = [
        "Live screens, not mockups",
        "Role-based onboarding",
        "NeoDove-ready operating flow",
        "Built for adoption, not just documentation",
    ]
    for idx, line in enumerate(summary_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(12)
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = BRAND_INK

    image_files = ["dashboard-today.png", "publish-content.png", "neodove-ops.png"]
    positions = [(Inches(7.25), Inches(1.35)), (Inches(9.05), Inches(2.1)), (Inches(7.75), Inches(3.95))]
    sizes = [(Inches(2.7), Inches(1.5)), (Inches(2.7), Inches(1.5)), (Inches(3.2), Inches(1.85))]
    for file_name, (left, top), (width, height) in zip(image_files, positions, sizes):
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left - Inches(0.05), top - Inches(0.05), width + Inches(0.1), height + Inches(0.1))
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = WHITE
        image_path = SCREENSHOT_DIR / file_name
        draw_w, draw_h = image_fit(image_path, width, height)
        draw_left = left + (width - draw_w) / 2
        draw_top = top + (height - draw_h) / 2
        slide.shapes.add_picture(str(image_path), draw_left, draw_top, width=draw_w, height=draw_h)

    add_textbox(slide, Inches(0.75), Inches(6.75), Inches(4.5), Inches(0.3), "Santaan Growth OS", 11, RGBColor(219, 227, 225))


def add_content_slide(prs: Presentation, spec: DeckSpec, index: int, section: DeckSection) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_rect(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height), BRAND_CREAM)
    fill_rect(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.35), prs.slide_height), BRAND_TEAL)

    fill_rect(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.55), Inches(1.75), Inches(0.42)), BRAND_AMBER)
    add_textbox(slide, Inches(0.92), Inches(0.62), Inches(1.4), Inches(0.2), f"SLIDE {index + 1}", 10, WHITE, bold=True)

    add_textbox(slide, Inches(0.7), Inches(1.1), Inches(5.4), Inches(0.7), section.title, 24, BRAND_INK, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.72), Inches(5.6), Inches(0.6), section.subtitle, 12, BRAND_MUTED)

    add_bullet_block(slide, Inches(0.7), Inches(2.25), Inches(4.0), Inches(4.6), section.bullets, font_size=15)

    if section.image_name:
        image_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.2), Inches(7.65), Inches(5.85))
        image_panel.fill.solid()
        image_panel.fill.fore_color.rgb = WHITE
        image_panel.line.color.rgb = BRAND_STONE
        image_path = SCREENSHOT_DIR / section.image_name
        frame_left = Inches(5.18)
        frame_top = Inches(1.38)
        frame_w = Inches(7.28)
        frame_h = Inches(5.48)
        draw_w, draw_h = image_fit(image_path, frame_w, frame_h)
        draw_left = frame_left + (frame_w - draw_w) / 2
        draw_top = frame_top + (frame_h - draw_h) / 2
        slide.shapes.add_picture(str(image_path), draw_left, draw_top, width=draw_w, height=draw_h)
    else:
        callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.0), Inches(7.65), Inches(3.5))
        fill_rect(callout, BRAND_TEAL)
        tf = callout.text_frame
        tf.clear()
        for idx, bullet in enumerate(section.bullets[:2]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(14)
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = WHITE

    add_textbox(slide, Inches(0.7), Inches(6.83), Inches(6.0), Inches(0.25), spec.title, 10, BRAND_MUTED)


def draw_wrapped(canvas_obj: canvas.Canvas, text: str, x: float, y: float, width: float, font_name: str, font_size: int, color_hex: str) -> float:
    styles = getSampleStyleSheet()
    style = ParagraphStyle(
        "wrap",
        parent=styles["BodyText"],
        fontName=font_name,
        fontSize=font_size,
        leading=font_size * 1.35,
        textColor=colors.HexColor(color_hex),
    )
    para = Paragraph(text, style)
    _, height = para.wrap(width, 1000)
    para.drawOn(canvas_obj, x, y - height)
    return y - height


def draw_pdf_image(c: canvas.Canvas, image_path: Path, x: float, y_top: float, max_w: float, max_h: float) -> None:
    with Image.open(image_path) as img:
        width, height = img.size
    scale = min(max_w / width, max_h / height)
    draw_w = width * scale
    draw_h = height * scale
    draw_x = x + (max_w - draw_w) / 2
    draw_y = y_top - draw_h - (max_h - draw_h) / 2
    c.drawImage(ImageReader(str(image_path)), draw_x, draw_y, width=draw_w, height=draw_h, preserveAspectRatio=True, mask="auto")


def build_pptx(spec: DeckSpec) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover_slide(prs, spec)
    for idx, section in enumerate(spec.sections):
        add_content_slide(prs, spec, idx, section)

    prs.save(str(spec.pptx_out))


def build_pdf(spec: DeckSpec) -> None:
    page_w, page_h = landscape(A4)
    c = canvas.Canvas(str(spec.pdf_out), pagesize=(page_w, page_h))

    # cover
    c.setFillColor(colors.HexColor("#225552"))
    c.rect(0, 0, page_w, page_h, fill=1, stroke=0)
    c.setFillColor(colors.HexColor("#e6a44e"))
    c.roundRect(34, page_h - 58, 120, 24, 8, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(48, page_h - 44, spec.kicker)
    c.setFont("Helvetica-Bold", 26)
    c.drawString(34, page_h - 98, spec.title)
    c.setFont("Helvetica", 14)
    c.setFillColor(colors.HexColor("#dbe3e1"))
    c.drawString(34, page_h - 122, spec.subtitle)

    c.setFillColor(colors.HexColor("#f6f2eb"))
    c.roundRect(34, page_h - 298, 315, 150, 14, fill=1, stroke=0)
    c.setFillColor(colors.HexColor("#1f262f"))
    c.setFont("Helvetica-Bold", 16)
    summary_lines = [
        "Live screens, not mockups",
        "Role-based onboarding",
        "NeoDove-ready operating flow",
        "Built for adoption, not just documentation",
    ]
    current_y = page_h - 178
    for line in summary_lines:
        c.drawString(56, current_y, line)
        current_y -= 27

    image_specs = [
        ("dashboard-today.png", 410, page_h - 150, 150, 90),
        ("publish-content.png", 535, page_h - 240, 150, 90),
        ("neodove-ops.png", 455, page_h - 368, 190, 115),
    ]
    for file_name, x, y, w, h in image_specs:
        c.setFillColor(colors.white)
        c.roundRect(x - 3, y - h - 3, w + 6, h + 6, 10, fill=1, stroke=0)
        draw_pdf_image(c, SCREENSHOT_DIR / file_name, x, y, w, h)
    c.showPage()

    # content pages
    for idx, section in enumerate(spec.sections, start=2):
        c.setFillColor(colors.HexColor("#f6f2eb"))
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        c.setFillColor(colors.HexColor("#225552"))
        c.rect(0, 0, 18, page_h, fill=1, stroke=0)

        c.setFillColor(colors.HexColor("#e6a44e"))
        c.roundRect(36, page_h - 54, 82, 20, 8, fill=1, stroke=0)
        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 10)
        c.drawString(52, page_h - 41, f"SLIDE {idx}")

        c.setFillColor(colors.HexColor("#1f262f"))
        c.setFont("Helvetica-Bold", 22)
        c.drawString(36, page_h - 86, section.title)
        c.setFillColor(colors.HexColor("#656d77"))
        c.setFont("Helvetica", 12)
        c.drawString(36, page_h - 108, section.subtitle)

        c.setFillColor(colors.white)
        c.roundRect(36, 54, 238, page_h - 184, 14, fill=1, stroke=0)
        c.setStrokeColor(colors.HexColor("#e0dcd5"))
        c.roundRect(36, 54, 238, page_h - 184, 14, fill=0, stroke=1)

        current_y = page_h - 150
        for bullet in section.bullets:
            current_y = draw_wrapped(c, f"• {bullet}", 52, current_y, 205, "Helvetica", 12, "#1f262f") - 10

        if section.image_name:
            c.setFillColor(colors.white)
            c.roundRect(302, 54, page_w - 338, page_h - 148, 14, fill=1, stroke=0)
            c.setStrokeColor(colors.HexColor("#e0dcd5"))
            c.roundRect(302, 54, page_w - 338, page_h - 148, 14, fill=0, stroke=1)
            draw_pdf_image(c, SCREENSHOT_DIR / section.image_name, 318, page_h - 128, page_w - 370, page_h - 190)

        c.setFillColor(colors.HexColor("#656d77"))
        c.setFont("Helvetica", 9)
        footer = f"{spec.title} - Page {idx}"
        c.drawString(page_w - stringWidth(footer, "Helvetica", 9) - 28, 18, footer)
        c.showPage()

    c.save()


def main() -> None:
    ensure_dirs()
    for deck in DECKS:
        build_pptx(deck)
        build_pdf(deck)
        print(deck.pptx_out)
        print(deck.pdf_out)


if __name__ == "__main__":
    main()
